"""
PowerPoint document handler (.pptx format).

Handles PowerPoint presentations with:
- Text extraction from slides
- Text extraction from notes
- Text extraction from tables and shapes
- Presentation reconstruction
"""

from pathlib import Path
from typing import Dict, Any, List, Optional
from datetime import datetime
import io

from .base import BaseDocument, DocumentFormat, DocumentMetadata
from ..logging_config import get_logger
from ..exceptions import DocumentProcessingError

logger = get_logger(__name__)


class PowerPointDocument(BaseDocument):
    """
    PowerPoint document handler for .pptx files.

    Supports:
    - Text extraction from slide content
    - Text extraction from notes
    - Text extraction from tables and shapes
    - Metadata extraction
    - Presentation reconstruction with anonymized content

    Note: Requires python-pptx library.
    """

    def __init__(self, file_path: Path):
        """
        Initialize PowerPoint document handler.

        Args:
            file_path: Path to PowerPoint file

        Raises:
            DocumentProcessingError: If PowerPoint document can't be loaded
        """
        self._presentation = None
        self._slides_map = []  # Track slide structure

        super().__init__(file_path)

        # Load document
        self._load_document()

    def _load_document(self):
        """Load PowerPoint document with python-pptx"""
        try:
            from pptx import Presentation

            self._presentation = Presentation(self.file_path)

            logger.info(
                f"PowerPoint document loaded successfully",
                extra={
                    "extra_fields": {
                        "slides": len(self._presentation.slides),
                        "file_size": self.get_file_size(),
                    }
                },
            )

        except ImportError:
            raise DocumentProcessingError(
                "python-pptx not installed. Install with: pip install python-pptx",
                {"required_package": "python-pptx"},
            )
        except Exception as e:
            logger.error(f"Failed to load PowerPoint document: {e}", exc_info=True)
            raise DocumentProcessingError(
                f"Failed to load PowerPoint document: {str(e)}",
                {"file_path": str(self.file_path)}
            )

    def _detect_format(self) -> DocumentFormat:
        """Detect format (always POWERPOINT for this handler)"""
        return DocumentFormat.POWERPOINT

    def _extract_metadata(self) -> DocumentMetadata:
        """Extract PowerPoint document metadata"""
        file_stats = self.file_path.stat()

        # Extract document properties
        author = None
        title = None
        creation_date = None
        modification_date = None

        try:
            if hasattr(self._presentation, "core_properties"):
                props = self._presentation.core_properties
                author = props.author
                title = props.title
                creation_date = props.created
                modification_date = props.modified

        except Exception as e:
            logger.warning(f"Failed to extract PowerPoint metadata: {e}")

        # Count slides
        slide_count = len(self._presentation.slides) if self._presentation else 0

        return DocumentMetadata(
            file_name=self.file_path.name,
            file_size=file_stats.st_size,
            format=DocumentFormat.POWERPOINT,
            page_count=slide_count,  # Slides as "pages"
            author=author,
            title=title,
            creation_date=creation_date,
            modification_date=modification_date,
            is_scanned=False,  # PowerPoint files are digital
            custom={
                "slides": slide_count,
            }
        )

    def extract_text(self) -> str:
        """
        Extract all text from PowerPoint document.

        Extracts from:
        - Slide titles and content
        - Text in shapes
        - Tables
        - Notes

        Returns:
            Extracted text content

        Raises:
            DocumentProcessingError: If extraction fails
        """
        try:
            text_parts = []

            for slide_idx, slide in enumerate(self._presentation.slides, start=1):
                slide_texts = []

                logger.debug(f"Extracting text from slide {slide_idx}")

                # Extract text from all shapes
                for shape in slide.shapes:
                    # Text frames (titles, content, text boxes)
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)

                    # Tables
                    if shape.has_table:
                        table_texts = []
                        for row in shape.table.rows:
                            row_texts = [cell.text.strip() for cell in row.cells]
                            if any(row_texts):
                                table_texts.append(" | ".join(row_texts))

                        if table_texts:
                            slide_texts.append("\n".join(table_texts))

                # Add slide content
                if slide_texts:
                    slide_content = f"[Slide {slide_idx}]\n" + "\n".join(slide_texts)
                    text_parts.append(slide_content)

                    self._slides_map.append({
                        "slide_number": slide_idx,
                        "text_index": len(text_parts) - 1,
                        "content": "\n".join(slide_texts)
                    })

                # Extract notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        text_parts.append(f"[Slide {slide_idx} - Notes]\n{notes_text}")

            full_text = "\n\n".join(text_parts)

            logger.info(
                f"Text extraction completed",
                extra={
                    "extra_fields": {
                        "total_slides": len(self._presentation.slides),
                        "total_length": len(full_text),
                    }
                },
            )

            return full_text

        except Exception as e:
            logger.error(f"Text extraction failed: {e}", exc_info=True)
            raise DocumentProcessingError(
                f"Failed to extract text from PowerPoint document: {str(e)}",
                {"file_path": str(self.file_path)}
            )

    def rebuild(self, anonymized_text: str, detections: List[Dict[str, Any]]) -> bytes:
        """
        Rebuild PowerPoint document with anonymized content.

        Creates a new presentation with anonymized text.

        Args:
            anonymized_text: Anonymized text
            detections: List of detections (for reference)

        Returns:
            PowerPoint document bytes
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            logger.info("Rebuilding PowerPoint document with anonymized content")

            # Create new presentation
            new_pres = Presentation()

            # Parse anonymized text back into slides
            sections = anonymized_text.split("\n\n")

            for section in sections:
                if not section.strip():
                    continue

                # Check if it's a slide section
                if section.startswith("[Slide") and "Notes]" not in section:
                    # Extract slide content
                    lines = section.split("\n")
                    slide_header = lines[0]  # [Slide N]
                    content_lines = lines[1:]

                    # Create slide
                    slide_layout = new_pres.slide_layouts[1]  # Title and Content
                    slide = new_pres.slides.add_slide(slide_layout)

                    # Add title (first line)
                    if content_lines:
                        title = slide.shapes.title
                        title.text = content_lines[0]

                    # Add content (remaining lines)
                    if len(content_lines) > 1:
                        content = slide.placeholders[1]
                        content.text = "\n".join(content_lines[1:])

                elif "[Notes]" in section:
                    # Notes sections - add to previous slide
                    # (Skip for now - notes reconstruction is complex)
                    pass

            # Save to bytes
            buffer = io.BytesIO()
            new_pres.save(buffer)
            ppt_bytes = buffer.getvalue()
            buffer.close()

            logger.info(
                "PowerPoint document rebuilt successfully",
                extra={"extra_fields": {"output_size": len(ppt_bytes)}},
            )

            return ppt_bytes

        except Exception as e:
            logger.error(f"PowerPoint document rebuild failed: {e}", exc_info=True)
            raise DocumentProcessingError(
                f"Failed to rebuild PowerPoint document: {str(e)}",
                {"file_path": str(self.file_path)}
            )

    def get_slide_count(self) -> int:
        """
        Get number of slides.

        Returns:
            Slide count
        """
        if hasattr(self, "_presentation") and self._presentation:
            return len(self._presentation.slides)
        return 0

    def close(self):
        """Close document resources"""
        if hasattr(self, "_presentation"):
            self._presentation = None
            logger.debug("PowerPoint document closed")

    def __del__(self):
        """Cleanup on deletion"""
        self.close()
